"""Size-aware extractor (spec §5).

Decision matrix:
  size > 50 MB                        -> filename_only
  size > 20 MB OR (pdf AND pages>30)  -> head_only (first 5-10 pages)
  else                                -> full
"""
from __future__ import annotations

from dataclasses import dataclass
from pathlib import Path
from typing import Optional

TWENTY_MB = 20 * 1024 * 1024
FIFTY_MB = 50 * 1024 * 1024
LONG_PDF_PAGES = 30
HEAD_PAGES = 8


@dataclass(frozen=True)
class ExtractionResult:
    text: str
    doc_type: str
    page_count: Optional[int]
    ocr_pages: tuple[int, ...] = ()


def decide_mode(size: int, page_count: Optional[int]) -> str:
    if size > FIFTY_MB:
        return "filename_only"
    if size > TWENTY_MB or (page_count is not None and page_count > LONG_PDF_PAGES):
        return "head_only"
    return "full"


def _detect_type(path: Path) -> str:
    ext = path.suffix.lower().lstrip(".")
    return {
        "pdf": "pdf",
        "docx": "docx", "doc": "docx",
        "pptx": "pptx", "ppt": "pptx",
        "xlsx": "xlsx", "xls": "xlsx",
    }.get(ext, ext or "bin")


def extract_text(path: Path, mode: str) -> ExtractionResult:
    doc_type = _detect_type(path)
    if mode == "filename_only":
        return ExtractionResult(text=path.name, doc_type=doc_type, page_count=None)
    if doc_type == "pdf":
        return _extract_pdf(path, mode)
    if doc_type == "docx":
        return _extract_docx(path)
    if doc_type == "pptx":
        return _extract_pptx(path)
    if doc_type == "xlsx":
        return _extract_xlsx(path)
    return ExtractionResult(text="", doc_type=doc_type, page_count=None)


def _render_page_rgb(page, dpi: int = 200):
    import numpy as np
    pix = page.get_pixmap(dpi=dpi, alpha=False)
    img = np.frombuffer(pix.samples, dtype=np.uint8).reshape(pix.height, pix.width, pix.n)
    return img[..., :3] if pix.n >= 3 else img


def _extract_pdf(path: Path, mode: str) -> ExtractionResult:
    import fitz  # PyMuPDF
    doc = fitz.open(str(path))
    total = doc.page_count
    limit = HEAD_PAGES if mode == "head_only" else total
    parts: list[str] = []
    ocr_pages: list[int] = []
    ocr_fn = None
    for i in range(min(limit, total)):
        page = doc.load_page(i)
        txt = page.get_text()
        if len(txt.strip()) < 10:
            ocr_pages.append(i)
            if ocr_fn is None:
                try:
                    from colab.ocr import ocr_page as ocr_fn_impl
                    ocr_fn = ocr_fn_impl
                except Exception as e:
                    print(f"[extractor] OCR unavailable: {type(e).__name__}: {e}",
                          flush=True)
                    ocr_fn = False
            if ocr_fn:
                try:
                    img = _render_page_rgb(page)
                    ocr_txt = ocr_fn(img)
                    if ocr_txt.strip():
                        parts.append(ocr_txt)
                except Exception as e:
                    print(f"[extractor] OCR page {i} failed on {path.name}: "
                          f"{type(e).__name__}: {e}", flush=True)
            continue
        parts.append(txt)
    return ExtractionResult(
        text="\n\n".join(parts),
        doc_type="pdf",
        page_count=total,
        ocr_pages=tuple(ocr_pages),
    )


def _extract_docx(path: Path) -> ExtractionResult:
    import docx
    try:
        doc = docx.Document(str(path))
    except Exception as e:  # PackageNotFoundError, BadZipFile, ancien .doc binaire
        print(f"[extractor] SKIP docx {path.name}: {type(e).__name__}", flush=True)
        return ExtractionResult(text=path.name, doc_type="docx", page_count=None)
    parts = [p.text for p in doc.paragraphs if p.text.strip()]
    return ExtractionResult(text="\n\n".join(parts), doc_type="docx", page_count=None)


def _extract_pptx(path: Path) -> ExtractionResult:
    from pptx import Presentation
    try:
        prs = Presentation(str(path))
    except Exception as e:
        print(f"[extractor] SKIP pptx {path.name}: {type(e).__name__}", flush=True)
        return ExtractionResult(text=path.name, doc_type="pptx", page_count=None)
    parts: list[str] = []
    for slide in prs.slides:
        slide_parts: list[str] = []
        for shape in slide.shapes:
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    txt = "".join(run.text for run in para.runs).strip()
                    if txt:
                        slide_parts.append(txt)
        if slide_parts:
            parts.append("\n".join(slide_parts))
    return ExtractionResult(
        text="\n\n".join(parts), doc_type="pptx",
        page_count=len(prs.slides),
    )


def _extract_xlsx(path: Path) -> ExtractionResult:
    from openpyxl import load_workbook
    try:
        wb = load_workbook(str(path), read_only=True, data_only=True)
    except Exception as e:
        print(f"[extractor] SKIP xlsx {path.name}: {type(e).__name__}", flush=True)
        return ExtractionResult(text=path.name, doc_type="xlsx", page_count=None)
    parts: list[str] = []
    sheet_count = len(wb.worksheets)
    for ws in wb.worksheets:
        sheet_parts: list[str] = [f"# {ws.title}"]
        for row in ws.iter_rows(values_only=True):
            cells = [str(c) for c in row if c is not None and str(c).strip()]
            if cells:
                sheet_parts.append(" | ".join(cells))
        if len(sheet_parts) > 1:
            # Blank line between rows so split_paragraphs yields one paragraph per row,
            # allowing the chunker to group them into token-sized chunks instead of
            # accumulating an entire sheet into a single unsplittable paragraph.
            parts.append("\n\n".join(sheet_parts))
    wb.close()
    return ExtractionResult(
        text="\n\n".join(parts), doc_type="xlsx",
        page_count=sheet_count,
    )
